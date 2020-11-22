import pyttsx3
import PyPDF2
import docx
from pptx import Presentation
import glob

#------Extracting text file data -------

def text_file_reader(path):

    with open(path, 'rb') as txtfile:
        return txtfile.readlines()

#------Extracting PDF file data -------

def pdf_file_reader(path):

     file = open(path, "rb")
     reader = PyPDF2.PdfFileReader(file)
     pages = reader.numPages
     pdfData = ""
     for i in range(0, pages - 1):
         pdfPage = reader.getPage(i)
         pdfData = pdfPage.extractText()
         pdfData = pdfData + pdfData
     return pdfData

#------Extracting Power point file data -------

def ppt_file_reader(path):

    pptdata = ""
    for eachfile in glob.glob("*.pptx"):
        prs = Presentation(eachfile)
        print(eachfile)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    data = shape.text
                    pptdata = pptdata + data
        return pptdata

#------Extracting Word file data -------

def word_file_reader(path):

    doc = docx.Document(path)
    finaltext = ""
    all_paras = doc.paragraphs
    print(len(all_paras))
    for para in all_paras:
        paratext = para.text
        finaltext= finaltext + paratext
    return finaltext

#-----Read and speak specified file data ----------

def speaker(filepath):

    filedetails = filepath.split(".")

    if filedetails[-1] == 'docx':
        text_data = word_file_reader(filepath)
    elif filedetails[-1] == 'pdf':
        text_data  = pdf_file_reader(filepath)
    elif filedetails[-1] == 'txt':
        text_data = text_file_reader(filepath)
    elif filedetails[-1] == 'pptx':
        text_data = ppt_file_reader(filepath)
    else:
        text_data ="Sorry !! Could not given any valid file for reading"
    init_speaker = pyttsx3.init()
    init_speaker.say(text_data)
    init_speaker.runAndWait()

#-----------START SPEAKER------------------

#path = "E:\My_Book_Reader\pywinauto.pdf"
#path = "E:\My_Book_Reader\TheITQC.pptx"
#path = "E:\My_Book_Reader\TheITQC.docx"
path = "E:\My_Book_Reader\TheITQC.txt"
speaker(path)
#-------------------------------------------
